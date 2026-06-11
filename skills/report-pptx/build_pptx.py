"""reports/{종목명}.md 를 읽어 디자인된 PPTX 리포트(reports/{종목명}.pptx)를 생성한다.

사용:
    python .claude/skills/report-pptx/build_pptx.py "reports/삼성전자.md"
    python .claude/skills/report-pptx/build_pptx.py "reports/삼성전자.md" --ticker 005930

설계 원칙(고정):
- 슬라이드 순서: 표지 → 종목 개요 → 재무 요약 → 가격/추세(차트) → 뉴스·심리 → 리스크 → 한 줄 종합
- 포인트색 KB 옐로우(#FFBC00) + 그레이·화이트 본문, 차분한 금융 리포트 톤
- 한글 폰트는 '맑은 고딕' 하나로 고정(글자 깨짐 방지)
- 매수/매도·목표가 단정 표현은 입력에서 그대로 옮기되 생성하지 않음(이 스크립트는 새 수치/의견을 만들지 않는다)
- 표가 슬라이드 밖으로 넘치면 행 수를 잘라 안전하게 맞춘다
"""

import argparse
import re
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# 디자인 상수 (KB 톤)
# ---------------------------------------------------------------------------
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)   # 포인트색
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)   # 본문 진한 글자
GRAY_MID = RGBColor(0x66, 0x66, 0x66)    # 보조 텍스트
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)  # 표 줄무늬/배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"                         # 한글 폰트 고정

# 16:9 슬라이드 기준 (EMU). 1인치 = 914400 EMU
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN * 2

# 재무 표 최대 행 수(헤더 포함). 넘치면 잘라 슬라이드 밖 침범 방지.
MAX_TABLE_ROWS = 9


# ---------------------------------------------------------------------------
# 마크다운 파싱
# ---------------------------------------------------------------------------
def parse_markdown(md_text: str) -> dict:
    """마크다운을 H2(##) 헤딩 기준으로 섹션 딕셔너리로 나눈다.

    반환: {"_title": 첫 H1 텍스트, "_date": 본문에서 찾은 기준일,
           "sections": [(heading, body_text), ...]}
    """
    lines = md_text.splitlines()
    title = ""
    sections = []
    cur_head = None
    cur_body: list[str] = []

    for line in lines:
        h1 = re.match(r"^#\s+(.*)", line)
        h2 = re.match(r"^##\s+(.*)", line)
        if h1 and not title:
            title = h1.group(1).strip()
            continue
        if h2:
            if cur_head is not None:
                sections.append((cur_head, "\n".join(cur_body).strip()))
            cur_head = h2.group(1).strip()
            cur_body = []
        else:
            if cur_head is not None:
                cur_body.append(line)
    if cur_head is not None:
        sections.append((cur_head, "\n".join(cur_body).strip()))

    # 기준일/작성일 추출 (예: "기준일: 2026-06-11")
    m = re.search(r"(?:기준일|작성일)\s*[:：]\s*(\d{4}-\d{2}-\d{2})", md_text)
    report_date = m.group(1) if m else date.today().isoformat()

    return {"_title": title, "_date": report_date, "sections": sections}


def find_section(sections, keywords) -> str:
    """헤딩에 키워드 중 하나가 들어간 첫 섹션 본문을 반환. 없으면 빈 문자열."""
    for head, body in sections:
        for kw in keywords:
            if kw in head:
                return body
    return ""


def extract_table(body: str):
    """본문에서 첫 마크다운 표를 파싱해 (헤더리스트, 행리스트)로 반환. 없으면 None."""
    rows = []
    for line in body.splitlines():
        s = line.strip()
        if s.startswith("|") and s.endswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            # 구분선(---) 행은 건너뜀
            if all(re.fullmatch(r":?-{2,}:?", c or "-") for c in cells):
                continue
            rows.append(cells)
    if len(rows) < 2:
        return None
    header, data = rows[0], rows[1:]
    return header, data


def extract_ticker(md_text: str) -> str | None:
    """본문에서 6자리 종목코드를 찾는다 (예: '삼성전자(005930)')."""
    m = re.search(r"\b(\d{6})\b", md_text)
    return m.group(1) if m else None


def body_to_bullets(body: str, limit: int = 7) -> list[tuple[str, int]]:
    """표·코드블록을 제외한 본문을 (텍스트, 들여쓰기레벨) 불릿 목록으로 변환."""
    bullets = []
    in_code = False
    for line in body.splitlines():
        s = line.strip()
        if s.startswith("```"):
            in_code = not in_code
            continue
        if in_code or not s:
            continue
        if s.startswith("|"):          # 표 행은 별도 표로 처리하므로 제외
            continue
        # 마크다운 마커 제거
        m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)", line)
        if m:
            indent = len(m.group(1))
            level = 1 if indent >= 2 else 0
            text = m.group(3)
        elif s.startswith("###"):
            text = s.lstrip("#").strip()
            level = 0
        else:
            text = s
            level = 0
        text = _clean_inline(text)
        if text:
            bullets.append((text, level))
        if len(bullets) >= limit:
            break
    return bullets


def _clean_inline(text: str) -> str:
    """**굵게**, [링크](url) 등 인라인 마크다운을 평문으로 정리."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1", text)  # 링크 텍스트만 남김
    return text.strip()


# ---------------------------------------------------------------------------
# PPTX 빌딩 헬퍼
# ---------------------------------------------------------------------------
def _set_font(run, size, bold=False, color=GRAY_DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 동아시아 폰트도 동일하게 지정해 한글 깨짐 방지
    rpr = run._r.get_or_add_rPr()
    ea = rpr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea", {}
    )
    ea.set("typeface", FONT)
    rpr.append(ea)


def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = 빈 레이아웃


def _add_accent_bar(slide, top=Inches(1.15)):
    """제목 아래 KB 옐로우 강조 바."""
    bar = slide.shapes.add_shape(
        1, MARGIN, top, Inches(1.4), Pt(5)  # 1 = 사각형
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_YELLOW
    bar.line.fill.background()
    return bar


def _add_header(slide, title, kicker=None):
    """모든 내용 슬라이드 상단 공통 헤더(키커 + 제목 + 강조 바)."""
    if kicker:
        tf = _add_textbox(slide, MARGIN, Inches(0.35), CONTENT_W, Inches(0.3))
        p = tf.paragraphs[0]
        _set_font(p.add_run(), 12, bold=True, color=GRAY_MID)
        p.runs[0].text = kicker
    tf = _add_textbox(slide, MARGIN, Inches(0.6), CONTENT_W, Inches(0.6))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_font(r, 26, bold=True, color=GRAY_DARK)
    _add_accent_bar(slide)


def _add_footer(slide, text="투자 판단은 사람 · 본 자료는 정보 제공 목적"):
    tf = _add_textbox(slide, MARGIN, SLIDE_H - Inches(0.45), CONTENT_W, Inches(0.3))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _set_font(r, 9, color=GRAY_MID)


# ---------------------------------------------------------------------------
# 슬라이드 7종
# ---------------------------------------------------------------------------
def slide_cover(prs, name, report_date):
    slide = _blank_slide(prs)
    # 배경: 진회색 좌측 띠 + 흰 본문
    band = slide.shapes.add_shape(1, 0, 0, Inches(0.35), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = KB_YELLOW
    band.line.fill.background()

    tf = _add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(1.2))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = name
    _set_font(r, 44, bold=True, color=GRAY_DARK)

    tf2 = _add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11), Inches(0.6))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "종목 리서치 리포트"
    _set_font(r2, 20, color=GRAY_MID)

    # 가드레일 고지: 첫머리 "무료 공개 데이터 기반 학습용"
    tfd = _add_textbox(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(0.4))
    rd = tfd.paragraphs[0].add_run()
    rd.text = "무료 공개 데이터 기반 학습용"
    _set_font(rd, 13, bold=True, color=KB_YELLOW)

    tf3 = _add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.5))
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = f"작성일 {report_date}"
    _set_font(r3, 14, color=GRAY_MID)
    return slide


def slide_bullets(prs, title, kicker, body, bullet_limit=7):
    slide = _blank_slide(prs)
    _add_header(slide, title, kicker)
    bullets = body_to_bullets(body, limit=bullet_limit)
    tf = _add_textbox(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(5.2))
    if not bullets:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "내용 없음 (입력 .md 에 해당 섹션이 비어 있음)"
        _set_font(r, 14, color=GRAY_MID)
    else:
        for i, (text, level) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            # 레벨0엔 KB 옐로우 불릿 점, 레벨1엔 들여쓰기
            mark = p.add_run()
            mark.text = ("●  " if level == 0 else "–  ")
            _set_font(mark, 14 if level == 0 else 12,
                      color=KB_YELLOW if level == 0 else GRAY_MID)
            r = p.add_run()
            r.text = text
            _set_font(r, 15 if level == 0 else 13,
                      bold=(level == 0), color=GRAY_DARK)
            p.level = level
    _add_footer(slide)
    return slide


def slide_financials(prs, body):
    slide = _blank_slide(prs)
    _add_header(slide, "재무 요약", "FINANCIALS · 최근 3개년")
    table = extract_table(body)
    top = Inches(1.6)

    if table:
        header, data = table
        # 행 넘침 방지: 헤더 + 데이터 최대 (MAX_TABLE_ROWS-1)
        max_data = MAX_TABLE_ROWS - 1
        trimmed = len(data) > max_data
        data = data[:max_data]
        ncols = len(header)
        nrows = len(data) + 1

        tbl_w = CONTENT_W
        tbl_h = Inches(0.45) * nrows
        gfx = slide.shapes.add_table(nrows, ncols, MARGIN, top, tbl_w, tbl_h)
        tbl = gfx.table
        # 기본 스타일 제거 후 직접 색칠
        for c in range(ncols):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = KB_YELLOW
            _cell_text(cell, header[c], 12, bold=True, color=GRAY_DARK,
                       align=PP_ALIGN.CENTER)
        for ri, row in enumerate(data, start=1):
            for ci in range(ncols):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else GRAY_LIGHT
                val = row[ci] if ci < len(row) else ""
                align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
                _cell_text(cell, _clean_inline(val), 11, color=GRAY_DARK,
                           align=align)
        if trimmed:
            note = _add_textbox(slide, MARGIN, top + tbl_h + Inches(0.1),
                                CONTENT_W, Inches(0.3))
            r = note.paragraphs[0].add_run()
            r.text = "※ 표 행이 많아 일부만 표시했습니다."
            _set_font(r, 9, color=GRAY_MID)
        # 표 아래 출처/코멘트 불릿
        cmt_top = top + tbl_h + Inches(0.4)
    else:
        cmt_top = top

    bullets = [b for b in body_to_bullets(body, limit=4)]
    if bullets:
        tf = _add_textbox(slide, MARGIN, cmt_top, CONTENT_W, Inches(2.0))
        for i, (text, level) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            mark = p.add_run()
            mark.text = "●  "
            _set_font(mark, 12, color=KB_YELLOW)
            r = p.add_run()
            r.text = text
            _set_font(r, 12, color=GRAY_DARK)
    _add_footer(slide)
    return slide


def _cell_text(cell, text, size, bold=False, color=GRAY_DARK, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_font(r, size, bold=bold, color=color)


def slide_price(prs, body, ticker, report_date, out_dir):
    slide = _blank_slide(prs)
    _add_header(slide, "가격 · 추세", "PRICE & TREND")
    chart_path = None
    if ticker:
        chart_path = _make_price_chart(ticker, out_dir, report_date)

    if chart_path and chart_path.exists():
        # 차트는 좌측, 코멘트는 우측 (이미지는 pptx 안에 복사되므로 임시파일은 곧 삭제)
        slide.shapes.add_picture(str(chart_path), MARGIN, Inches(1.55),
                                 width=Inches(7.4))
        try:
            chart_path.unlink()
        except OSError:
            pass
        tf = _add_textbox(slide, Inches(8.3), Inches(1.55),
                          Inches(4.4), Inches(5.0))
        bullets = body_to_bullets(body, limit=6)
        for i, (text, level) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            mark = p.add_run()
            mark.text = "●  "
            _set_font(mark, 12, color=KB_YELLOW)
            r = p.add_run()
            r.text = text
            _set_font(r, 12, color=GRAY_DARK)
    else:
        # 차트 실패 시 텍스트 슬라이드로 폴백
        bullets = body_to_bullets(body, limit=8)
        tf = _add_textbox(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(5.0))
        if ticker:
            note = tf.paragraphs[0].add_run()
            note.text = "(차트 생성 실패 — 텍스트 요약으로 대체)"
            _set_font(note, 10, color=GRAY_MID)
        for text, level in bullets:
            p = tf.add_paragraph()
            mark = p.add_run()
            mark.text = "●  "
            _set_font(mark, 13, color=KB_YELLOW)
            r = p.add_run()
            r.text = text
            _set_font(r, 14, color=GRAY_DARK)
    _add_footer(slide)
    return slide


def _make_price_chart(ticker, out_dir, report_date):
    """FinanceDataReader로 최근 6개월 종가+20/60일선 차트 PNG 생성."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import font_manager
        import FinanceDataReader as fdr
        from datetime import datetime, timedelta

        # 한글 폰트
        malgun = r"C:\Windows\Fonts\malgun.ttf"
        if Path(malgun).exists():
            fp = font_manager.FontProperties(fname=malgun)
            plt.rcParams["font.family"] = fp.get_name()
        plt.rcParams["axes.unicode_minus"] = False

        end = datetime.fromisoformat(report_date)
        start = end - timedelta(days=200)
        df = fdr.DataReader(ticker, start.strftime("%Y-%m-%d"),
                            end.strftime("%Y-%m-%d"))
        if df is None or df.empty:
            return None
        df = df.tail(130)
        ma20 = df["Close"].rolling(20).mean()
        ma60 = df["Close"].rolling(60).mean()

        fig, ax = plt.subplots(figsize=(7.4, 4.6), dpi=150)
        ax.plot(df.index, df["Close"], color="#333333", lw=1.6, label="종가")
        ax.plot(df.index, ma20, color="#FFBC00", lw=1.4, label="20일선")
        ax.plot(df.index, ma60, color="#999999", lw=1.2, ls="--", label="60일선")
        ax.set_title(f"{ticker} 최근 6개월 주가 (출처: FinanceDataReader, {report_date})",
                     fontsize=11, color="#333333")
        ax.legend(loc="upper left", fontsize=9, frameon=False)
        ax.grid(True, color="#EEEEEE")
        for spine in ["top", "right"]:
            ax.spines[spine].set_visible(False)
        ax.tick_params(labelsize=8, colors="#666666")
        fig.tight_layout()

        out_dir.mkdir(parents=True, exist_ok=True)
        png = out_dir / f"_chart_{ticker}.png"
        fig.savefig(png, bbox_inches="tight")
        plt.close(fig)
        return png
    except Exception as e:  # 차트 실패해도 리포트 생성은 계속
        print(f"[chart] 생성 실패: {e}", file=sys.stderr)
        return None


def slide_oneline(prs, body):
    slide = _blank_slide(prs)
    # 한 줄 종합은 강조 슬라이드: 옐로우 배경 박스 + 큰 텍스트
    _add_header(slide, "한 줄 종합", "SUMMARY")
    box = slide.shapes.add_shape(1, MARGIN, Inches(2.3), CONTENT_W, Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LIGHT
    box.line.color.rgb = KB_YELLOW
    box.line.width = Pt(2)

    # 본문에서 가장 핵심적인 한두 줄 추출
    bullets = body_to_bullets(body, limit=3)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    if bullets:
        for i, (text, level) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = text
            _set_font(r, 18, bold=(i == 0), color=GRAY_DARK)
    else:
        r = tf.paragraphs[0].add_run()
        r.text = "종합 의견 없음"
        _set_font(r, 18, color=GRAY_MID)
    _add_footer(slide, "투자 판단은 사람 · 매수/매도·목표가 의견을 제시하지 않습니다")
    return slide


def slide_sources(prs, body):
    """가드레일: 리포트 끝에 데이터 출처·기준일 목록."""
    slide = _blank_slide(prs)
    _add_header(slide, "데이터 출처 · 기준일", "SOURCES")
    bullets = body_to_bullets(body, limit=12)
    tf = _add_textbox(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(4.6))
    if bullets:
        for i, (text, level) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            mark = p.add_run()
            mark.text = "●  "
            _set_font(mark, 12, color=KB_YELLOW)
            r = p.add_run()
            r.text = text
            _set_font(r, 12, color=GRAY_DARK)
    else:
        r = tf.paragraphs[0].add_run()
        r.text = "출처 목록이 입력 .md 에 없습니다 (가드레일: 출처·기준일 목록 필수)."
        _set_font(r, 13, color=GRAY_MID)
    _add_footer(slide, "무료 공개 데이터 기반 학습용 · 투자 판단은 사람")
    return slide


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------
def build(md_path: Path, ticker: str | None = None) -> Path:
    md_text = md_path.read_text(encoding="utf-8")
    parsed = parse_markdown(md_text)
    name = md_path.stem  # 파일명 = 종목명
    report_date = parsed["_date"]
    sections = parsed["sections"]
    ticker = ticker or extract_ticker(md_text)
    out_dir = md_path.parent

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 순서 고정
    slide_cover(prs, name, report_date)
    slide_bullets(prs, "종목 개요", "OVERVIEW",
                  find_section(sections, ["개요", "Overview", "소개"]))
    slide_financials(prs, find_section(sections, ["재무", "Financial", "실적"]))
    slide_price(prs, find_section(sections, ["가격", "추세", "기술", "Price"]),
                ticker, report_date, out_dir)
    slide_bullets(prs, "뉴스 · 심리", "NEWS & SENTIMENT",
                  find_section(sections, ["뉴스", "심리", "Sentiment", "이슈"]))
    slide_bullets(prs, "리스크", "RISK",
                  find_section(sections, ["리스크", "Risk", "위험"]))
    slide_oneline(prs, find_section(sections, ["한 줄", "종합", "결론", "Summary"]))
    slide_sources(prs, find_section(sections, ["출처", "Source", "데이터 출처"]))

    out_path = out_dir / f"{name}.pptx"
    prs.save(str(out_path))
    return out_path


def main():
    ap = argparse.ArgumentParser(description="종목 리서치 .md → 디자인 PPTX")
    ap.add_argument("md", help="입력 마크다운 경로 (reports/{종목명}.md)")
    ap.add_argument("--ticker", help="6자리 종목코드(차트용). 미지정 시 본문에서 탐색.")
    args = ap.parse_args()

    md_path = Path(args.md)
    if not md_path.exists():
        sys.exit(f"입력 파일 없음: {md_path}")
    out = build(md_path, args.ticker)
    print(f"생성 완료: {out}")


if __name__ == "__main__":
    main()
